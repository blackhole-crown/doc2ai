"""PPTX 文件读取器 - 带 OCR 支持（改进临时文件清理）"""

import os
import tempfile
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None


class PptxOcrReader:
    """PPTX 文件读取器 - 支持 OCR 提取图片文字"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.metadata = {
            'slides': 0,
            'images_extracted': 0,
            'ocr_success': 0,
            'ocr_failed': 0,
            'error': None,
            'has_text': False
        }
        self.use_ocr = config.get('use_ocr', True)
        self.ocr_lang = config.get('ocr_lang', 'chi_sim+eng')
        
        # 配置 Tesseract 路径
        tesseract_cmd = config.get('tesseract_cmd')
        if tesseract_cmd and pytesseract:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    
    def _ocr_image(self, img_path):
        """对单张图片进行 OCR，确保临时文件清理"""
        temp_img = None
        try:
            if not self.use_ocr or not pytesseract or not Image:
                return None
            
            img = Image.open(img_path)
            
            # 简单预处理
            if img.mode != 'L':
                img = img.convert('L')
            
            # 使用 pytesseract 识别
            text = pytesseract.image_to_string(img, lang=self.ocr_lang)
            
            return text.strip()
            
        except Exception as e:
            print(f"    OCR 失败: {e}")
            return None
    
    def read(self, file_path):
        """读取 PPTX 文件，提取文本和图片 OCR"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            self.metadata['error'] = 'file_not_found'
            return None
        
        if Presentation is None:
            self.metadata['error'] = 'python_pptx_not_installed'
            print("错误: python-pptx 未安装，请运行: pip install python-pptx")
            return None
        
        if self.use_ocr and (Image is None or pytesseract is None):
            print("警告: OCR 库未安装，将只提取文本")
            self.use_ocr = False
        
        text_parts = []
        ocr_success = 0
        ocr_failed = 0
        temp_dir = None
        
        try:
            # 打开 PPT
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)
            print(f"PPT 总页数: {total_slides}")
            
            # 创建临时目录
            if self.use_ocr:
                temp_dir = tempfile.mkdtemp()
            
            # 提取文本和图片
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # 1. 提取形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                    
                    # 2. 尝试提取图片并 OCR
                    if self.use_ocr and hasattr(shape, "image") and shape.image:
                        try:
                            image_blob = shape.image.blob
                            if image_blob:
                                img_temp = os.path.join(temp_dir, f"slide_{slide_num}_img_{hash(image_blob)}.png")
                                with open(img_temp, 'wb') as f:
                                    f.write(image_blob)
                                
                                # OCR 识别
                                text = self._ocr_image(img_temp)
                                if text and text.strip():
                                    slide_texts.append(f"[图片识别] {text.strip()}")
                                    ocr_success += 1
                                else:
                                    ocr_failed += 1
                                
                                # 立即删除临时文件
                                if os.path.exists(img_temp):
                                    os.remove(img_temp)
                                    
                        except Exception as e:
                            ocr_failed += 1
                
                # 记录这一页
                if slide_texts:
                    text_parts.append(f"=== 第 {slide_num} 页 ===")
                    text_parts.extend(slide_texts)
                    text_parts.append("")
            
            if not text_parts:
                self.metadata['error'] = 'no_content'
                print(f"⚠️  PPT 文件没有内容")
                return None
            
            full_text = '\n'.join(text_parts)
            
            # 更新元数据
            self.metadata.update({
                'slides': total_slides,
                'ocr_success': ocr_success,
                'ocr_failed': ocr_failed,
                'has_text': True
            })
            
            print(f"✅ PPT 处理完成")
            print(f"   OCR 成功: {ocr_success} 张")
            print(f"   OCR 失败: {ocr_failed} 张")
            
            return full_text
            
        except Exception as e:
            self.metadata['error'] = str(e)
            print(f"读取 PPT 失败: {e}")
            import traceback
            traceback.print_exc()
            return None
        finally:
            # 清理临时目录
            if temp_dir and os.path.exists(temp_dir):
                import shutil
                shutil.rmtree(temp_dir, ignore_errors=True)
    
    def get_metadata(self):
        return self.metadata