"""PPTX 文件读取器 - 增强版"""

import os
from pathlib import Path
import re

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("警告: python-pptx 未安装，请运行: pip install python-pptx")

# 可选：尝试导入 python-pptx 的高级功能
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    MSO_SHAPE_TYPE = None


class PptxReader:
    """PowerPoint 文档读取器 - 增强版"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.metadata = {
            'slides': 0,
            'truncated': False,
            'line_count': 0,
            'displayed_lines': 0,
            'size': 0,
            'error': None,
            'has_text': False
        }
    
    def _extract_shape_text(self, shape):
        """提取形状中的文本，包括分组形状中的文本"""
        texts = []
        
        # 如果有文本属性
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                texts.append(text)
        
        # 如果是分组形状，递归提取子形状
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                texts.extend(self._extract_shape_text(sub_shape))
        
        # 如果是表格形状
        if hasattr(shape, "table"):
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))
        
        return texts
    
    def read(self, file_path):
        """读取 PPTX 文件"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            self.metadata['error'] = 'file_not_found'
            print(f"文件不存在: {file_path}")
            return None
        
        if Presentation is None:
            self.metadata['error'] = 'python_pptx_not_installed'
            print("错误: python-pptx 未安装，请运行: pip install python-pptx")
            return None
        
        try:
            # 打开演示文稿
            prs = Presentation(str(file_path))
            
            text_parts = []
            total_slides = len(prs.slides)
            slide_has_text = []
            
            print(f"PPT 总页数: {total_slides}")
            
            # 提取每一页的内容
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # 1. 提取标题
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_texts.append(f"【标题】{title}")
                
                # 2. 遍历所有形状提取文本
                for shape in slide.shapes:
                    # 跳过标题（已处理）
                    if shape == slide.shapes.title:
                        continue
                    
                    # 提取文本
                    shape_texts = self._extract_shape_text(shape)
                    if shape_texts:
                        slide_texts.extend(shape_texts)
                
                # 3. 提取备注（如果有）
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_texts.append(f"【备注】{notes_text}")
                
                # 记录这一页是否有文本
                if slide_texts:
                    slide_has_text.append(True)
                    text_parts.append(f"=== 第 {slide_num} 页 ===")
                    text_parts.extend(slide_texts)
                    text_parts.append("")
                else:
                    slide_has_text.append(False)
            
            # 如果所有页面都没有文本
            if not text_parts:
                self.metadata['error'] = 'no_text_content'
                self.metadata['has_text'] = False
                print(f"⚠️  PPT 文件没有文本内容（可能是图片形式）: {file_path}")
                print(f"   建议：尝试使用 OCR 工具提取图片中的文字")
                return None
            
            full_text = '\n'.join(text_parts)
            self.metadata['has_text'] = True
            
            # 显示统计
            text_slides = sum(slide_has_text)
            print(f"   有文本的页面: {text_slides}/{total_slides}")
            
            # 限制行数
            max_lines = self.config.get('max_lines', 500)
            lines = full_text.split('\n')
            original_lines = len(lines)
            
            if len(lines) > max_lines:
                lines = lines[:max_lines//2] + ['... 省略中间内容 ...'] + lines[-max_lines//2:]
                full_text = '\n'.join(lines)
                truncated = True
            else:
                truncated = False
            
            # 更新元数据
            self.metadata.update({
                'slides': total_slides,
                'slide_count': total_slides,
                'text_slides': text_slides,
                'truncated': truncated,
                'line_count': original_lines,
                'displayed_lines': len(lines),
                'size': file_path.stat().st_size,
                'error': None
            })
            
            print(f"✅ PPT 提取完成，共 {len(full_text)} 字符，{text_slides}/{total_slides} 页有文本")
            
            return full_text
            
        except Exception as e:
            self.metadata['error'] = str(e)
            print(f"读取 PPT 失败: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def get_metadata(self):
        """返回元数据"""
        return self.metadata


class PptxReaderWithOCR:
    """带 OCR 支持的 PPT 读取器（需要额外安装 pytesseract）"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.metadata = {}
        self.use_ocr = config.get('use_ocr', False)
    
    def read(self, file_path):
        """读取 PPTX 文件，可选 OCR"""
        # 先尝试普通提取
        reader = PptxReader(self.config)
        content = reader.read(file_path)
        
        # 如果没有文本且启用了 OCR，尝试提取图片
        if not content and self.use_ocr:
            print("尝试使用 OCR 提取图片中的文字...")
            try:
                from PIL import Image
                import pytesseract
                from pptx import Presentation
                
                prs = Presentation(str(file_path))
                text_parts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    # 这里需要将幻灯片转换为图片，然后 OCR
                    # 这需要额外的库：python-pptx 本身不支持直接转图片
                    # 可以考虑使用 pywin32（Windows）或 LibreOffice 命令行
                    print(f"  第 {slide_num} 页 OCR 功能待实现")
                
                if text_parts:
                    content = '\n'.join(text_parts)
                    print(f"OCR 提取完成")
                    return content
            except Exception as e:
                print(f"OCR 失败: {e}")
        
        return content
    
    def get_metadata(self):
        return self.metadata
    def read_by_xml(self, file_path):
        """使用 XML 直接提取文本（备用方法）"""
        try:
            import zipfile
            import re
            
            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                slide_files = [f for f in zip_ref.namelist() 
                            if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                slide_files.sort()
                
                for i, slide_file in enumerate(slide_files, 1):
                    with zip_ref.open(slide_file) as f:
                        content = f.read().decode('utf-8')
                        texts = re.findall(r'<a:t>(.*?)</a:t>', content)
                        
                        if texts:
                            text_parts.append(f"=== 第 {i} 页 ===")
                            text_parts.extend([t for t in texts if t.strip()])
                            text_parts.append("")
            
            if text_parts:
                return '\n'.join(text_parts)
        except Exception as e:
            print(f"XML 提取失败: {e}")
        
        return None